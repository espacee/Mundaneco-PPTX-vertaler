#!/usr/bin/env python
import os
import sys
import time
from pptx import Presentation
import openai
from dotenv import load_dotenv

# Load environment variables from .env and set the API key.
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
if not openai.api_key:
    print("Error: OPENAI_API_KEY is not set in the environment.")
    sys.exit(1)

def load_presentation(file_path):
    """
    Loads a .pptx file and returns a Presentation object.
    """
    try:
        presentation = Presentation(file_path)
        return presentation
    except Exception as e:
        print(f"Error loading file {file_path}: {e}")
        sys.exit(1)

def translate_text(text, target_language):
    """
    Translates the given text from Dutch to the specified target_language 
    using OpenAI's GPT-4-turbo model.
    """
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "system", 
                    "content": (
                        f"You are a professional translator. Please translate the following Dutch text "
                        f"into {target_language}, preserving its original formatting where possible."
                    )
                },
                {"role": "user", "content": f"Text:\n{text}"}
            ],
            temperature=0
        )
        translated_text = response.choices[0].message['content'].strip()
        return translated_text
    except Exception as e:
        print(f"Translation error: {e}")
        return text  # Fallback to original text in case of error.

def translate_presentation(input_file, output_file, target_language):
    """
    Translates the text on each slide from the input presentation and saves 
    the new presentation with the translated content.
    """
    presentation = load_presentation(input_file)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not shape.has_text_frame:
                continue
            original_text = shape.text_frame.text
            if original_text.strip() == "":
                continue
            print(f"Translating slide {slide_index}, shape {shape_index}...")
            translated_text = translate_text(original_text, target_language)
            # Update the text frame with the translated text.
            shape.text_frame.text = translated_text
            # Optional: Wait briefly to help avoid rate-limit issues.
            time.sleep(1)
    presentation.save(output_file)

def main():
    if len(sys.argv) < 2:
        print("Usage: python translate_pptx.py <input.pptx>")
        sys.exit(1)

    input_file = sys.argv[1]
    if not os.path.exists(input_file):
        print(f"Error: File {input_file} does not exist.")
        sys.exit(1)

    file_base, file_ext = os.path.splitext(input_file)
    output_file_fr = f"{file_base}_fr{file_ext}"
    output_file_en = f"{file_base}_en{file_ext}"

    print("Starting translation to French...")
    translate_presentation(input_file, output_file_fr, "French")
    print(f"French translation saved to {output_file_fr}")

    print("Starting translation to English...")
    translate_presentation(input_file, output_file_en, "English")
    print(f"English translation saved to {output_file_en}")

if __name__ == "__main__":
    main() 